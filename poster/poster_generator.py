"""
Generate APS 2026 Job Seeker Poster as PowerPoint (.pptx).

Matches the Typst poster_v4.typ layout: 36x48 landscape, light theme,
3 project cards, skills section, contact footer.

Usage:
    python poster/poster_generator.py
    python poster/poster_generator.py --output poster/poster.pptx
"""

import argparse
import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logging.basicConfig(level=logging.INFO, format='%(asctime)s %(message)s')
logger = logging.getLogger(__name__)

ROOT = Path(__file__).resolve().parent.parent
FIGURES_DIR = ROOT / 'paper' / 'figures'

# Colors (Cream & Navy palette)
PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)       # Deep navy
SECONDARY = RGBColor(0x2C, 0x3E, 0x5A)     # Medium navy
ACCENT_SLATE = RGBColor(0x3D, 0x50, 0x68)  # Slate navy
TEXT_PRIMARY = RGBColor(0x1A, 0x18, 0x14)   # Warm near-black
TEXT_SECONDARY = RGBColor(0x5C, 0x56, 0x4E) # Warm dark gray
TEXT_MUTED = RGBColor(0x8A, 0x83, 0x7A)     # Warm medium gray
CREAM = RGBColor(0xFA, 0xF7, 0xF2)         # Warm cream background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # Clean white (card bg)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)        # Clean white
CARD_BORDER = RGBColor(0xD6, 0xCF, 0xC5)   # Warm taupe
TAG_BG = RGBColor(0xED, 0xE8, 0xE1)        # Light warm tan
FOOTER_TEXT = RGBColor(0xFA, 0xF7, 0xF2)    # Cream on navy

# Poster dimensions: 48 x 36 inches
POSTER_W = Inches(48)
POSTER_H = Inches(36)
MARGIN_X = Inches(0.6)
MARGIN_Y = Inches(0.5)
CONTENT_W = Inches(48 - 1.2)


def add_textbox(slide, left, top, width, height, text, font_size=24,
                bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT,
                font_name='Helvetica Neue', italic=False):
    """Add a text box with single-run formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    run.font.italic = italic
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG,
                     line_color=CARD_BORDER, line_width=Pt(1.5)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = line_width
    return shape


def add_accent_bar(slide, left, top, width, color):
    """Add a colored accent bar (thin rectangle)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_image_or_placeholder(slide, left, top, width, height, image_path, label):
    """Add an image if it exists, otherwise a placeholder rectangle."""
    if image_path.exists():
        pic = slide.shapes.add_picture(str(image_path), left, top, width, height)
        return pic
    else:
        shape = add_rounded_rect(slide, left, top, width, height,
                                fill_color=CREAM)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(14)
        run.font.color.rgb = TEXT_MUTED
        return shape


def build_project_card(slide, left, top, width, height, title, accent_color,
                       hook, problem, approach_bullets, outcome_bullets,
                       tools, figure_path, figure_label):
    """Build a project card with title, content, figure, and tools."""
    # Card background
    card = add_rounded_rect(slide, left, top, width, height)

    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    pad = Inches(0.25)
    inner_left = left + Inches(0.25)
    inner_w = width - Inches(0.5)
    y = top + pad

    # Title
    add_textbox(slide, inner_left, y, inner_w, Inches(0.5),
               title, font_size=30, bold=True, color=accent_color)
    y += Inches(0.5)

    # Accent line
    add_accent_bar(slide, inner_left, y, inner_w, accent_color)
    y += Inches(0.15)

    # Hook
    add_textbox(slide, inner_left, y, inner_w, Inches(0.4),
               hook, font_size=17, italic=True, color=accent_color)
    y += Inches(0.45)

    # Problem
    add_textbox(slide, inner_left, y, inner_w, Inches(0.3),
               'Problem', font_size=20, bold=True, color=TEXT_PRIMARY)
    y += Inches(0.3)
    add_textbox(slide, inner_left, y, inner_w, Inches(0.5),
               problem, font_size=17, color=TEXT_PRIMARY)
    y += Inches(0.55)

    # Approach
    add_textbox(slide, inner_left, y, inner_w, Inches(0.3),
               'Approach', font_size=20, bold=True, color=TEXT_PRIMARY)
    y += Inches(0.3)
    for bullet in approach_bullets:
        tb = add_textbox(slide, inner_left + Inches(0.15), y, inner_w - Inches(0.15), Inches(0.4),
                        f"\u2022 {bullet}", font_size=16, color=TEXT_PRIMARY)
        y += Inches(0.35)

    y += Inches(0.1)

    # Figure
    fig_h = Inches(2.8)
    add_image_or_placeholder(slide, inner_left, y, inner_w, fig_h,
                            figure_path, figure_label)
    y += fig_h + Inches(0.15)

    # Outcome
    add_textbox(slide, inner_left, y, inner_w, Inches(0.3),
               'Outcome', font_size=20, bold=True, color=TEXT_PRIMARY)
    y += Inches(0.3)
    for bullet in outcome_bullets:
        add_textbox(slide, inner_left + Inches(0.15), y, inner_w - Inches(0.15), Inches(0.4),
                   f"\u2022 {bullet}", font_size=16, color=TEXT_PRIMARY)
        y += Inches(0.35)

    y += Inches(0.1)

    # Tools
    tools_text = '  |  '.join(tools)
    add_textbox(slide, inner_left, y, inner_w, Inches(0.3),
               tools_text, font_size=14, color=TEXT_SECONDARY)


def build_poster():
    """Build the complete APS Job Seeker Poster as PowerPoint."""
    prs = Presentation()
    prs.slide_width = POSTER_W
    prs.slide_height = POSTER_H

    # Blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Cream background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CREAM

    # ================================================================
    # HEADER
    # ================================================================
    y = MARGIN_Y

    # Headshot placeholder
    headshot = add_rounded_rect(
        slide, MARGIN_X, y, Inches(1.6), Inches(2.0),
        fill_color=CREAM, line_color=PRIMARY
    )
    headshot.line.width = Pt(3)
    tf = headshot.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'HEADSHOT'
    run.font.size = Pt(14)
    run.font.color.rgb = TEXT_MUTED
    run.font.bold = True

    # Name
    name_left = MARGIN_X + Inches(2.0)
    name_w = Inches(34)

    # "WILL" in blue, "HAMMOND" in purple — need two textboxes
    add_textbox(slide, name_left, y, Inches(6), Inches(0.8),
               'WILL', font_size=50, bold=True, color=PRIMARY)
    add_textbox(slide, name_left + Inches(5.2), y, Inches(12), Inches(0.8),
               'HAMMOND', font_size=50, bold=True, color=SECONDARY)

    y_sub = y + Inches(0.75)
    add_textbox(slide, name_left, y_sub, Inches(30), Inches(0.4),
               'Pitzer College  \u00B7  B.A. Physics & Mathematics  \u00B7  Expected 2028',
               font_size=22, color=TEXT_SECONDARY)

    y_contact = y_sub + Inches(0.4)
    add_textbox(slide, name_left, y_contact, Inches(30), Inches(0.3),
               'whammond@pitzer.edu  |  linkedin.com/in/willhammond  |  github.com/willhammondhimself',
               font_size=17, color=TEXT_MUTED)

    # Title quote
    y_title = y_contact + Inches(0.5)
    add_textbox(slide, name_left, y_title, Inches(34), Inches(0.6),
               '\u201cQuantitative Research at the Intersection of Physics, Math, and Finance\u201d',
               font_size=32, bold=True, color=SECONDARY, italic=True)

    # About Me
    y_about = y_title + Inches(0.7)
    about_text = (
        "I'm a sophomore studying physics and mathematics with a minor in data science, "
        "taking graduate-level coursework across the Claremont Colleges. I love applying "
        "mathematical tools to complex, uncertain systems \u2014 from detecting financial market "
        "crises using differential geometry to building systematic trading strategies. "
        "Advised by Prof. Trung Phan, Keck Science Department."
    )
    add_textbox(slide, name_left, y_about, Inches(34), Inches(1.0),
               about_text, font_size=18, color=TEXT_PRIMARY)

    # QR code placeholders (right side)
    qr_left = Inches(48 - 0.6 - 1.3)
    for i, (label, color) in enumerate([
        ('LinkedIn', PRIMARY), ('GitHub', SECONDARY), ('Portfolio', ACCENT_SLATE)
    ]):
        qr_top = y + Inches(i * 1.35)
        qr = add_rounded_rect(slide, qr_left, qr_top, Inches(1.2), Inches(1.2),
                              fill_color=CREAM, line_color=color)
        qr.line.width = Pt(1.5)
        tf = qr.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f'QR: {label}'
        run.font.size = Pt(11)
        run.font.color.rgb = TEXT_MUTED

    # Divider line
    y_divider = y_about + Inches(1.2)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, y_divider, CONTENT_W, Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CARD_BORDER
    line.line.fill.background()

    # ================================================================
    # THREE PROJECT CARDS
    # ================================================================
    cards_top = y_divider + Inches(0.2)
    card_w = Inches(15.2)
    card_h = Inches(18.5)
    card_gap = Inches(0.3)

    # Card 1: QCML Research
    build_project_card(
        slide,
        left=MARGIN_X,
        top=cards_top,
        width=card_w,
        height=card_h,
        title='QCML Research',
        accent_color=PRIMARY,
        hook='I applied physics math to find hidden patterns in financial data',
        problem='Can we detect market regime shifts without labeled crisis data?',
        approach_bullets=[
            'Embed time series into geometric space using spectral metric learning',
            'Extract 3 geometric observables measuring manifold deformation during stress',
            'Fully unsupervised \u2014 no crisis labels, no look-ahead',
        ],
        outcome_bullets=[
            'Multi-Lag Fidelity (d=1.44) beat Random Forest (d=1.13, p=0.002)',
            'Tested across 5 ETFs and 4 out-of-sample crises',
            'Solo-authored 34-page research paper (in revision)',
        ],
        tools=['Python', 'PyTorch', 'NumPy/SciPy', 'Optuna', 'LaTeX'],
        figure_path=FIGURES_DIR / 'poster_qcml_vs_rf.png',
        figure_label='QCML vs RF Comparison',
    )

    # Card 2: Quanta Ventures
    build_project_card(
        slide,
        left=MARGIN_X + card_w + card_gap,
        top=cards_top,
        width=card_w,
        height=card_h,
        title='Quanta Ventures',
        accent_color=SECONDARY,
        hook='Developed systematic equity strategies for a quantitative fund',
        problem='Building robust multi-strategy portfolios that survive regime changes',
        approach_bullets=[
            'Modular framework: multiple strategy sleeves with volatility targeting',
            'VIX/VVIX regime diagnostics for dynamic leverage',
            'Walk-forward validation with 30-day embargo and robustness suite',
        ],
        outcome_bullets=[
            'Sharpe 2.92 and Calmar 5.02 on out-of-sample data',
            'Passed full robustness test suite',
            'Built C++ Heston pricer for options strategy component',
        ],
        tools=['Python', 'scikit-learn', 'Optuna', 'PostgreSQL', 'C++'],
        figure_path=FIGURES_DIR / 'poster_quanta_metrics.png',
        figure_label='Quanta Performance Metrics',
    )

    # Card 3: Student Quant Fund
    build_project_card(
        slide,
        left=MARGIN_X + 2 * (card_w + card_gap),
        top=cards_top,
        width=card_w,
        height=card_h,
        title='Student Quant Fund',
        accent_color=ACCENT_SLATE,
        hook='Finding alpha and building community',
        problem='Giving underclassmen access to quantitative skills and a production trading platform',
        approach_bullets=[
            'Teaching momentum and mean-variance strategies',
            'Building production-ready platform for future cohorts',
            'Preparing for paper trading and capital deployment',
        ],
        outcome_bullets=[
            'Mentoring newer students in collaborative strategy development',
            'Building institutional knowledge across cohorts',
            'Platform designed for live paper trading',
        ],
        tools=['Python', 'Jupyter', 'Team Collaboration'],
        figure_path=FIGURES_DIR / 'poster_skills_radar.png',
        figure_label='Skills Radar Chart',
    )

    # ================================================================
    # SKILLS SECTION
    # ================================================================
    skills_top = cards_top + card_h + Inches(0.2)
    skills_h = Inches(4.5)

    skills_bg = add_rounded_rect(
        slide, MARGIN_X, skills_top, CONTENT_W, skills_h,
        fill_color=CARD_BG, line_color=CARD_BORDER
    )

    add_textbox(slide, MARGIN_X + Inches(0.3), skills_top + Inches(0.2),
               Inches(4), Inches(0.5), 'Skills',
               font_size=28, bold=True, color=PRIMARY)

    # Hard Skills
    hs_left = MARGIN_X + Inches(0.3)
    hs_top = skills_top + Inches(0.8)
    add_textbox(slide, hs_left, hs_top, Inches(8), Inches(0.35),
               'Hard Skills', font_size=21, bold=True, color=TEXT_PRIMARY)

    hs_y = hs_top + Inches(0.4)
    add_textbox(slide, hs_left, hs_y, Inches(22), Inches(0.3),
               'Tech:  Python (NumPy, Pandas, PyTorch)  |  SQL  |  Git/CI-CD  |  Docker  |  LaTeX  |  Mathematica',
               font_size=16, color=TEXT_SECONDARY)
    hs_y += Inches(0.35)
    add_textbox(slide, hs_left, hs_y, Inches(22), Inches(0.3),
               'CS:  Data Structures & Algorithms  |  ML/PyTorch  |  Numerical Optimization',
               font_size=16, color=TEXT_SECONDARY)
    hs_y += Inches(0.35)
    add_textbox(slide, hs_left, hs_y, Inches(22), Inches(0.3),
               'Finance:  Stochastic Calculus  |  Options Pricing (BS, Heston)  |  Backtesting/Walk-Forward  |  Risk Metrics (Sharpe, VaR, CVaR)',
               font_size=16, color=TEXT_SECONDARY)

    # Soft Skills
    ss_left = MARGIN_X + Inches(23)
    ss_top = skills_top + Inches(0.8)
    add_textbox(slide, ss_left, ss_top, Inches(20), Inches(0.35),
               'Soft Skills', font_size=21, bold=True, color=TEXT_PRIMARY)

    soft_skills = [
        ('CS Teaching Assistant', 'Held weekly office hours, mentored 60+ students'),
        ('Independent Researcher', 'Solo-authored 34-page research paper'),
        ('Team Collaboration', 'Student Quant Fund leadership'),
        ('Technical Communication', 'Translating complex math for diverse audiences'),
    ]
    ss_y = ss_top + Inches(0.4)
    for title, desc in soft_skills:
        add_textbox(slide, ss_left, ss_y, Inches(20), Inches(0.25),
                   title, font_size=17, bold=True, color=TEXT_PRIMARY)
        ss_y += Inches(0.28)
        add_textbox(slide, ss_left, ss_y, Inches(20), Inches(0.25),
                   desc, font_size=15, color=TEXT_SECONDARY)
        ss_y += Inches(0.35)

    # ================================================================
    # FOOTER
    # ================================================================
    footer_top = skills_top + skills_h + Inches(0.15)
    footer_h = Inches(1.2)

    footer_bg = add_rounded_rect(
        slide, MARGIN_X, footer_top, CONTENT_W, footer_h,
        fill_color=PRIMARY, line_color=PRIMARY
    )

    # Left: institution + honors
    add_textbox(slide, MARGIN_X + Inches(0.3), footer_top + Inches(0.15),
               Inches(12), Inches(0.35),
               'Pitzer College  |  Keck Science Department',
               font_size=17, bold=True, color=FOOTER_TEXT)
    add_textbox(slide, MARGIN_X + Inches(0.3), footer_top + Inches(0.5),
               Inches(12), Inches(0.3),
               'Carpe Diem Endowed Scholar',
               font_size=15, color=RGBColor(0xB0, 0xA9, 0x9F))

    # Center: seeking
    add_textbox(slide, Inches(14), footer_top + Inches(0.15),
               Inches(22), Inches(0.4),
               'Seeking: Quant Research/Trading  |  ML/AI Engineering  |  Physics Research',
               font_size=21, bold=True, color=FOOTER_TEXT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(14), footer_top + Inches(0.55),
               Inches(22), Inches(0.3),
               'Internships & Full-time \u2014 Available May 2028',
               font_size=18, color=RGBColor(0xD6, 0xCF, 0xC5), alignment=PP_ALIGN.CENTER)

    # Right: interests
    add_textbox(slide, Inches(38), footer_top + Inches(0.25),
               Inches(9), Inches(0.5),
               'USA Climbing Divisionals  \u00B7  Poker  \u00B7  Guitar & Drums',
               font_size=16, color=RGBColor(0xB0, 0xA9, 0x9F), alignment=PP_ALIGN.RIGHT)

    return prs


def main():
    parser = argparse.ArgumentParser(description='Generate APS poster as PowerPoint')
    parser.add_argument('--output', type=str, default=None,
                       help='Output path (default: poster/poster.pptx)')
    args = parser.parse_args()

    output = Path(args.output) if args.output else ROOT / 'poster' / 'poster.pptx'

    logger.info("Building APS Job Seeker Poster (PowerPoint)...")
    prs = build_poster()

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    logger.info(f"Saved: {output}")
    logger.info(f"Dimensions: 48x36 inches")
    logger.info("Open in PowerPoint/Keynote to verify layout.")


if __name__ == '__main__':
    main()

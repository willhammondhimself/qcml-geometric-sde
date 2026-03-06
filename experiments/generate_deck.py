#!/usr/bin/env python3
"""Generate a PowerPoint deck summarizing the QCML Geometric Observatory research.

Reads pre-generated figures from paper/figures/ and produces a ~15-slide
presentation at paper/qcml_geometric_observatory_deck.pptx.

Usage:
    python experiments/generate_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
REPO = Path(__file__).resolve().parent.parent
FIGURES = REPO / "paper" / "figures"
OUTPUT = REPO / "paper" / "qcml_geometric_observatory_deck.pptx"

# ---------------------------------------------------------------------------
# Design tokens (navy + white, matching plot_style.py)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x2D, 0x6A, 0x6A)
BURGUNDY = RGBColor(0x8B, 0x22, 0x52)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
SLATE = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF6, 0xF0)

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"


def new_prs() -> Presentation:
    """Create a blank 16:9 presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def _set_bg(slide, color: RGBColor):
    """Set solid background fill on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  bold=False, color=NAVY, alignment=PP_ALIGN.LEFT,
                  font_name=FONT_BODY, line_spacing=1.2):
    """Add a text box to a slide and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * (line_spacing - 1))
    return tf


def _add_paragraph(tf, text, font_size=18, bold=False, color=NAVY,
                   alignment=PP_ALIGN.LEFT, bullet=False, font_name=FONT_BODY):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if bullet:
        p.level = 1
    return p


def add_title_bar(slide, title_text: str, subtitle: str = None):
    """Add a navy bar across the top with white title text."""
    # Navy rectangle across top
    left, top, width, height = Inches(0), Inches(0), SLIDE_W, Inches(1.4)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # Title text
    _add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11.5), Inches(0.7),
                  title_text, font_size=28, bold=True, color=WHITE, font_name=FONT_TITLE)

    if subtitle:
        _add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.45),
                      subtitle, font_size=16, color=GOLD, font_name=FONT_BODY)


def add_text_slide(prs, title: str, bullets: list[str], subtitle: str = None):
    """Add a slide with title bar and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, WHITE)
    add_title_bar(slide, title, subtitle)

    tf = _add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.2),
                       bullets[0], font_size=22, color=NAVY)
    for b in bullets[1:]:
        _add_paragraph(tf, b, font_size=22, color=NAVY)
    return slide


def add_figure_slide(prs, title: str, figure_name: str, caption: str = None,
                     subtitle: str = None):
    """Add a slide with title bar, centered figure, and optional caption."""
    fig_path = FIGURES / figure_name
    if not fig_path.exists():
        print(f"  WARNING: {fig_path} not found, adding placeholder text")
        return add_text_slide(prs, title, [f"[Figure missing: {figure_name}]"], subtitle)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, WHITE)
    add_title_bar(slide, title, subtitle)

    # Figure — fill available space below title bar with margin
    img_left = Inches(0.8)
    img_top = Inches(1.7)
    img_max_w = Inches(11.7)
    img_max_h = Inches(5.0) if not caption else Inches(4.5)

    slide.shapes.add_picture(str(fig_path), img_left, img_top,
                             width=img_max_w)

    if caption:
        _add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
                      caption, font_size=14, color=SLATE,
                      alignment=PP_ALIGN.CENTER)
    return slide


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------

def build_deck():
    prs = new_prs()

    # ------------------------------------------------------------------
    # 1. Title Slide
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, NAVY)

    _add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.5),
                  "The Geometric Observatory",
                  font_size=44, bold=True, color=WHITE, font_name=FONT_TITLE,
                  alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(1.0),
                  "for Financial Crisis Detection",
                  font_size=36, bold=False, color=GOLD, font_name=FONT_TITLE,
                  alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.8),
                  "22 instruments  |  17 crises  |  27 years of data",
                  font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

    # ------------------------------------------------------------------
    # 2. The Question
    # ------------------------------------------------------------------
    add_text_slide(prs, "The Question", [
        "Can you detect financial crises early by measuring",
        "the geometry of market returns?",
        "",
        "We built 22 geometric observables inspired by quantum",
        "information theory and tested them alongside 14 traditional",
        "methods across 17 historical crises spanning 27 years.",
    ])

    # ------------------------------------------------------------------
    # 3. The Observatory Analogy
    # ------------------------------------------------------------------
    slide = add_text_slide(prs, "The Observatory Analogy", [
        "A single weather sensor catches one type of event.",
        "A modern weather station has many instruments.",
        "",
        "Our 22 geometric channels span five families:",
    ])
    tf = slide.shapes[-1].text_frame  # last added text box
    families = [
        "Metric  (QFI, geodesic velocity, metric condition)",
        "Curvature  (Berry phase, sectional curvature, Ricci scalar)",
        "Spectral  (spectral gap, entropy, complexity)",
        "Topological  (Chern number, dimensional collapse)",
        "Information  (purity, fidelity, Hamiltonian sensitivity)",
    ]
    for f in families:
        _add_paragraph(tf, f, font_size=18, color=TEAL, bullet=True)

    # ------------------------------------------------------------------
    # 4. What We Tested — crisis timeline figure
    # ------------------------------------------------------------------
    add_figure_slide(prs, "What We Tested", "crisis_timeline_17.png",
                     caption="17 crises from 1997 (Asian Crisis) to 2024 (Carry Unwind) "
                             "— 36 total methods: 22 geometric + 14 baselines")

    # ------------------------------------------------------------------
    # 5. Finding #1 — Leaderboard
    # ------------------------------------------------------------------
    add_figure_slide(prs, "Finding 1: The Geometric Methods Work",
                     "ranked_methods_barchart.png",
                     subtitle="Reduced Purity #1 (d = 0.834)  |  8 of top 10 are geometric  "
                              "|  Random Forest dropped from #1 to #17")

    # ------------------------------------------------------------------
    # 6. Finding #1 Detail — Top 10 consistency
    # ------------------------------------------------------------------
    add_figure_slide(prs, "How Consistent Are the Top Methods?",
                     "effect_sizes_top10.png",
                     caption="Violin plots: distribution of Cohen's d across 17 crises. "
                             "Wide = inconsistent. Narrow = reliable.")

    # ------------------------------------------------------------------
    # 7. Finding #2 — Heatmap: no single winner
    # ------------------------------------------------------------------
    add_figure_slide(prs, "Finding 2: No Single Winner",
                     "crisis_heatmap_36x17.png",
                     subtitle="14 different methods win across 17 crises")

    # ------------------------------------------------------------------
    # 8. Finding #2 Detail — why specialization matters
    # ------------------------------------------------------------------
    add_text_slide(prs, "Each Crisis Has a Different Geometric Signature", [
        "2008 GFC  \u2192  VIX Level wins (d = 4.5) \u2014 obvious vol spike",
        "2018 Volmageddon  \u2192  Hamiltonian Sensitivity (d = 1.7) \u2014 subtle structural shift",
        "2021 Meme/Archegos  \u2192  Dim. Collapse (d = 1.7) \u2014 hidden stress under calm surface",
        "2016 Brexit  \u2192  Speed Limit Ratio (d = 1.2) \u2014 rapid state transition",
        "2022 Rate Hikes  \u2192  Reduced Purity (d = 1.3) \u2014 prolonged regime change",
        "",
        "Each geometric channel is like a different wavelength of light \u2014",
        "some crises are 'visible' only in certain channels.",
    ])

    # ------------------------------------------------------------------
    # 9. Finding #3 — Fusion beats individuals
    # ------------------------------------------------------------------
    add_figure_slide(prs, "Finding 3: Fusion Beats Individuals",
                     "fusion_comparison.png",
                     subtitle="Regime-Adaptive fusion (d = 0.774) "
                              "\u2014 top fusion, less variance than any individual champion")

    # ------------------------------------------------------------------
    # 10. Finding #3 Detail — Generalization test
    # ------------------------------------------------------------------
    add_figure_slide(prs, "The Generalization Test",
                     "holdout_generalization.png",
                     subtitle="Train on 15 crises, hold out 4 post-2020 crises")

    slide = prs.slides[-1]
    _add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7),
                  "Regime-Adaptive +1.1%  |  Reduced Purity \u221259%  |  "
                  "The combination is more robust than any part",
                  font_size=15, color=BURGUNDY, alignment=PP_ALIGN.CENTER)

    # ------------------------------------------------------------------
    # 11. Lead Times
    # ------------------------------------------------------------------
    add_figure_slide(prs, "Early Warning: Lead Times",
                     "lead_time_analysis.png",
                     caption="Many geometric methods detect crises 150\u2013200+ days in advance "
                             "\u2014 not just classification, actual early warning")

    # ------------------------------------------------------------------
    # 12. Narrative: GFC 2008
    # ------------------------------------------------------------------
    add_figure_slide(prs, "What It Looks Like: 2008 Global Financial Crisis",
                     "narrative_2008_gfc.png",
                     caption="8 detectors showing real-time signals through GFC. "
                             "Gold band = crisis period.")

    # ------------------------------------------------------------------
    # 13. Narrative: COVID 2020
    # ------------------------------------------------------------------
    add_figure_slide(prs, "What It Looks Like: 2020 COVID Crash",
                     "narrative_2020_covid.png",
                     caption="Same 8 detectors, different crisis \u2014 "
                             "different channels activate for different failure modes.")

    # ------------------------------------------------------------------
    # 14. Implications
    # ------------------------------------------------------------------
    add_text_slide(prs, "Implications", [
        "For finance:",
        "  Exploitable geometric structure in markets that vol-based methods miss entirely.",
        "",
        "For the field:",
        "  The 'one model to rule them all' approach is wrong for crisis detection.",
        "  Adaptive fusion across channels is the way forward.",
        "",
        "For quantum-inspired methods:",
        "  Hilbert space representation produces observables that are",
        "  both interpretable and empirically useful.",
    ])

    # ------------------------------------------------------------------
    # 15. Honest Caveats
    # ------------------------------------------------------------------
    add_text_slide(prs, "Honest Caveats", [
        "Median QCML \u2248 median baselines (0.326 vs 0.352).",
        "  The top geometric methods are exceptional; the average is comparable.",
        "",
        "Fusion Friedman test is non-significant (p = 0.538).",
        "  No fusion method statistically dominates. The case is about",
        "  consistency and generalization, not raw dominance.",
        "",
        "Lead times are noisy \u2014 some early detections may be false positives.",
        "",
        "Historical data only \u2014 no live forward test yet.",
    ])

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved deck to {OUTPUT}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    build_deck()
